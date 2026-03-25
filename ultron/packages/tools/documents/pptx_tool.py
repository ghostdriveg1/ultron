import os
import tempfile
from pptx import Presentation

from packages.tools.base_tool import BaseTool
from packages.tools.schemas.pptx_schema import PPTXInput, PPTXOutput
from .fast_io_client import FastIOClient

class PPTXTool(BaseTool):
    """Tool to create PowerPoint presentations."""
    input_schema = PPTXInput
    output_schema = PPTXOutput

    def __init__(self):
        super().__init__(
            name="create_pptx",
            description="Creates a PowerPoint presentation from a list of slides.",
            permission_level="ALWAYS_ALLOWED"
        )
        self.fast_io = FastIOClient()

    async def execute(self, params: PPTXInput) -> PPTXOutput:
        with tempfile.TemporaryDirectory() as temp_dir:
            file_path = os.path.join(temp_dir, params.output_filename)
            
            # Create PPTX
            prs = Presentation()
            
            # Use slide layout 1 (Title and Content)
            bullet_slide_layout = prs.slide_layouts[1]
            
            for slide_data in params.slides:
                slide = prs.slides.add_slide(bullet_slide_layout)
                shapes = slide.shapes
                
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                
                title_shape.text = slide_data.title
                
                tf = body_shape.text_frame
                tf.text = slide_data.content
                
            prs.save(file_path)
            
            # Upload
            fast_io_url = await self.fast_io.upload(file_path, params.output_filename)
            
            # Mock GitHub URL
            github_url = f"https://github.com/mock/outputs/{params.output_filename}"
            
            return PPTXOutput(
                fast_io_url=fast_io_url,
                github_url=github_url,
                slide_count=len(params.slides)
            )
